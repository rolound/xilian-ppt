from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os
import tempfile

app = Flask(__name__)

def create_pptx(title, slides_content):
    """生成PPT的核心逻辑，和之前一样，只是返回值改一下"""
    prs = Presentation()
    
    # 封面
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "由昔涟的记忆编织而成 ✨"
    
    # 内容页
    for i, slide_data in enumerate(slides_content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data.get("title", f"第{i+2}页")
        
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()
        for point in slide_data.get("content", []):
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    
    # 保存到临时文件
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(tmp.name)
    return tmp.name

@app.route('/create_pptx', methods=['POST'])
def handle_create_pptx():
    """这是服务员的前台，收到请求就干活"""
    try:
        data = request.get_json()
        title = data.get("title", "未命名PPT")
        slides_content = data.get("slides", [])
        
        filepath = create_pptx(title, slides_content)
        return send_file(
            filepath,
            as_attachment=True,
            download_name='xilian_pptx.pptx',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return jsonify({"error": f"昔涟的记忆有些混乱：{str(e)}"}), 500

